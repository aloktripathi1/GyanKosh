import io
import re

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from app.agents.document_type_hints import SCANNED_PDF
from app.llm.client import ocr_page_text
from app.schemas.document_intelligence import (
    DocumentIntelligenceOutput,
    DocumentSection,
    EquationBlock,
    FigureBlock,
    TableBlock,
)

EQUATION_PATTERN = re.compile(r"[=≈≠≤≥∑∫√±×÷^]|\b[a-zA-Z]\s*=\s*[-\w]")

# Below this many characters of extracted text, a page is almost certainly a
# scan or image-only — not a legitimately short page of real content.
OCR_FALLBACK_THRESHOLD_CHARS = 40
# A user-declared "scanned" hint trusts the signal more: even a page with a
# little stray extracted text (stamps, watermarks, running headers) still
# gets the OCR treatment rather than being taken at face value.
OCR_FALLBACK_THRESHOLD_CHARS_HINTED = 200


class DocumentIntelligenceInput:
    """Raw bytes, not a storage_path — the caller reads the file via the storage
    interface and hands this agent plain bytes, so parsing stays a pure function."""

    def __init__(self, document_id: str, file_type: str, content: bytes, document_type_hint: str | None = None):
        self.document_id = document_id
        self.file_type = file_type
        self.content = content
        self.document_type_hint = document_type_hint


def _find_equations(section_id: str, text: str, page: int | None) -> list[EquationBlock]:
    equations = []
    for i, line in enumerate(text.splitlines()):
        line = line.strip()
        if line and EQUATION_PATTERN.search(line) and len(line) < 200:
            equations.append(EquationBlock(id=f"{section_id}-eq{i}", raw_text=line, page=page))
    return equations


def _parse_pdf(document_id: str, content: bytes, document_type_hint: str | None = None) -> DocumentIntelligenceOutput:
    doc = fitz.open(stream=content, filetype="pdf")
    sections: list[DocumentSection] = []
    raw_text_parts: list[str] = []
    threshold = OCR_FALLBACK_THRESHOLD_CHARS_HINTED if document_type_hint == SCANNED_PDF else OCR_FALLBACK_THRESHOLD_CHARS

    for page_index, page in enumerate(doc):
        page_num = page_index + 1
        text = page.get_text()

        if len(text.strip()) < threshold:
            pixmap = page.get_pixmap(dpi=150)
            ocr_text = ocr_page_text(pixmap.tobytes("png"), media_type="image/png")
            if ocr_text and ocr_text != "[blank page]":
                text = ocr_text

        raw_text_parts.append(text)
        section_id = f"page-{page_num}"

        tables = []
        try:
            for t_index, table in enumerate(page.find_tables().tables):
                rows = table.extract()
                tables.append(TableBlock(id=f"{section_id}-table{t_index}", rows=[[str(c) if c is not None else "" for c in row] for row in rows], page=page_num))
        except Exception:
            pass

        figures = [
            FigureBlock(id=f"{section_id}-fig{i}", page=page_num)
            for i in range(len(page.get_images()))
        ]

        sections.append(
            DocumentSection(
                id=section_id,
                heading=None,
                level=1,
                text=text,
                page=page_num,
                tables=tables,
                figures=figures,
                equations=_find_equations(section_id, text, page_num),
            )
        )

    return DocumentIntelligenceOutput(
        document_id=document_id,
        title=doc.metadata.get("title") or None,
        sections=sections,
        page_count=doc.page_count,
        raw_text="\n".join(raw_text_parts),
    )


def _parse_docx(document_id: str, content: bytes) -> DocumentIntelligenceOutput:
    doc = DocxDocument(io.BytesIO(content))
    sections: list[DocumentSection] = []
    raw_text_parts: list[str] = []

    current_heading = None
    current_level = 1
    current_text: list[str] = []
    section_index = 0

    def flush():
        nonlocal section_index
        if current_text or current_heading:
            section_id = f"section-{section_index}"
            text = "\n".join(current_text)
            sections.append(
                DocumentSection(
                    id=section_id,
                    heading=current_heading,
                    level=current_level,
                    text=text,
                    equations=_find_equations(section_id, text, None),
                )
            )
            section_index += 1

    for para in doc.paragraphs:
        style = (para.style.name or "").lower()
        if style.startswith("heading"):
            flush()
            current_text = []
            current_heading = para.text
            current_level = int(style.replace("heading", "").strip() or 1)
        else:
            if para.text.strip():
                current_text.append(para.text)
                raw_text_parts.append(para.text)
    flush()

    for t_index, table in enumerate(doc.tables):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        sections.append(
            DocumentSection(
                id=f"table-{t_index}",
                heading=None,
                level=1,
                text="",
                tables=[TableBlock(id=f"table-{t_index}", rows=rows)],
            )
        )

    return DocumentIntelligenceOutput(
        document_id=document_id,
        title=doc.core_properties.title or None,
        sections=sections,
        page_count=None,
        raw_text="\n".join(raw_text_parts),
    )


def _parse_pptx(document_id: str, content: bytes) -> DocumentIntelligenceOutput:
    prs = Presentation(io.BytesIO(content))
    sections: list[DocumentSection] = []
    raw_text_parts: list[str] = []

    for slide_index, slide in enumerate(prs.slides):
        section_id = f"slide-{slide_index + 1}"
        heading = None
        texts: list[str] = []
        tables: list[TableBlock] = []
        figures: list[FigureBlock] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                if shape == slide.shapes.title:
                    heading = shape_text
                elif shape_text:
                    texts.append(shape_text)
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                tables.append(TableBlock(id=f"{section_id}-table{len(tables)}", rows=rows))
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                figures.append(FigureBlock(id=f"{section_id}-fig{len(figures)}"))

        text = "\n".join(texts)
        raw_text_parts.append(f"{heading or ''}\n{text}")
        sections.append(
            DocumentSection(
                id=section_id,
                heading=heading,
                level=1,
                text=text,
                tables=tables,
                figures=figures,
                equations=_find_equations(section_id, text, None),
            )
        )

    return DocumentIntelligenceOutput(
        document_id=document_id,
        title=None,
        sections=sections,
        page_count=len(prs.slides),
        raw_text="\n".join(raw_text_parts),
    )


def _parse_txt(document_id: str, content: bytes) -> DocumentIntelligenceOutput:
    text = content.decode("utf-8", errors="replace")
    section_id = "section-0"
    return DocumentIntelligenceOutput(
        document_id=document_id,
        title=None,
        sections=[
            DocumentSection(
                id=section_id,
                heading=None,
                level=1,
                text=text,
                equations=_find_equations(section_id, text, None),
            )
        ],
        page_count=None,
        raw_text=text,
    )


_PARSERS = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "txt": _parse_txt,
}


def run(input: DocumentIntelligenceInput) -> DocumentIntelligenceOutput:
    parser = _PARSERS.get(input.file_type)
    if parser is None:
        raise ValueError(f"Unsupported file_type for Document Intelligence: {input.file_type}")
    if input.file_type == "pdf":
        return _parse_pdf(input.document_id, input.content, input.document_type_hint)
    return parser(input.document_id, input.content)
